from io import BytesIO
import pandas as pd
from docx import Document
from pptx import Presentation
import os
import fitz  # PyMuPDF for PDFs
from fastapi import UploadFile
from starlette.datastructures import UploadFile as StarletteUploadFile
from pdf2docx import Converter  # Import pdf2docx
import markdown  # To parse Markdown files

async def convert_file_to_docx(file: UploadFile) -> UploadFile:
    _, ext = os.path.splitext(file.filename.lower())

    try:
        content = await file.read()
        buffer = BytesIO(content)

        # Save the uploaded file temporarily (for PDFs)
        temp_pdf_path = "temp_input.pdf"
        with open(temp_pdf_path, "wb") as f:
            f.write(content)

        # Create a Word document
        doc = Document()

        if ext == ".txt":
            doc.add_paragraph(content.decode("utf-8"))

        elif ext == ".csv":
            df = pd.read_csv(buffer)
            for row in df.itertuples(index=False):
                doc.add_paragraph("\t".join(map(str, row)))

        elif ext == ".xlsx":
            df = pd.read_excel(buffer)
            for row in df.itertuples(index=False):
                doc.add_paragraph("\t".join(map(str, row)))

        elif ext == ".docx":
            existing_doc = Document(buffer)
            for para in existing_doc.paragraphs:
                doc.add_paragraph(para.text)

        elif ext == ".pdf":
            output_path = "converted.docx"

            # Convert PDF to DOCX using pdf2docx
            cv = Converter(temp_pdf_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()

            # Load the saved DOCX file into memory buffer
            with open(output_path, "rb") as f:
                output_buffer = BytesIO(f.read())

            upload_file = StarletteUploadFile(filename="converted.docx", file=output_buffer)
            
            # Clean up temp file
            os.remove(temp_pdf_path)

            return upload_file

        elif ext == ".pptx":
            prs = Presentation(buffer)
            for i, slide in enumerate(prs.slides):
                doc.add_paragraph(f"Slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        doc.add_paragraph(shape.text)

        elif ext == ".html":
            doc.add_paragraph(content.decode("utf-8"))

        elif ext == ".md":
            md_content = content.decode("utf-8")
            html_content = markdown.markdown(md_content)  # Convert MD to HTML
            doc.add_paragraph(html_content)  # Add Markdown as plain text

        else:
            doc.add_paragraph("[Unsupported file type for conversion]")

        # Save the Word document
        output_path = "converted.docx"
        doc.save(output_path)

        # Load the saved file into memory buffer
        with open(output_path, "rb") as f:
            output_buffer = BytesIO(f.read())

        upload_file = StarletteUploadFile(filename="converted.docx", file=output_buffer)

    except Exception as e:
        error_msg = f"Error converting file: {str(e)}"
        error_buffer = BytesIO(error_msg.encode("utf-8"))
        upload_file = StarletteUploadFile(filename="error.docx", file=error_buffer)

    return upload_file
